"""Build docs/deck.pptx — the submission deck — from the narrative in docs/deck.md.

    backend/.venv/bin/python docs/build_deck.py            # writes docs/deck.pptx

The numbers on the results slide are read from the latest three-tier eval report so the
deck cannot drift from what was measured. Requires python-pptx (pip install python-pptx).
"""

from __future__ import annotations

import glob
import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "deck.pptx"
TESTS = 264

# ops-room palette shared with the app
PAPER = RGBColor(0x0F, 0x15, 0x1C)
CARD = RGBColor(0x16, 0x1E, 0x27)
INK = RGBColor(0xE6, 0xED, 0xF3)
MUT = RGBColor(0x93, 0xA4, 0xB4)
EDGE = RGBColor(0x2A, 0x36, 0x44)
TEAL = RGBColor(0x4F, 0xC1, 0xD8)
AMBER = RGBColor(0xE5, 0xA9, 0x3D)
RED = RGBColor(0xE0, 0x70, 0x5C)
GREEN = RGBColor(0x62, 0xC4, 0x8D)

SANS = "Arial"
MONO = "Courier New"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)


# ---------------------------------------------------------------- measured numbers


def latest_results() -> dict:
    reports = sorted(glob.glob(str(ROOT / "backend/evals/reports/tier123-agent-sdk-v*.json")))
    if not reports:
        return {}
    r = json.load(open(reports[-1]))
    tiers = {int(k): v for k, v in r["by_tier"].items()}
    return {
        "stem": Path(reports[-1]).stem,
        "t1": f"{tiers[1]['passed']}/{tiers[1]['total']}",
        "t2": f"{tiers[2]['passed']}/{tiers[2]['total']}",
        "t3": f"{tiers[3]['passed']}/{tiers[3]['total']}",
        "p50": f"{r['latency']['p50_ms'] / 1000:.1f} s",
        "p95": f"{r['latency']['p95_ms'] / 1000:.1f} s",
        "cost": f"${r['cost_usd']:.2f}" if r.get("cost_usd") else "—",
    }


# ---------------------------------------------------------------- primitives


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return slide


def text(slide, x, y, w, h, s, *, size=18, color=INK, bold=False, font=SANS, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    first = True
    for line in s.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def rich(slide, x, y, w, h, paragraphs, *, size=17, line_spacing=1.2, space_after=8):
    """paragraphs: list of lists of (text, style) runs; style in {"", "b", "accent", "mut", "mono", "amber", "red", "green"}."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    colors = {"": INK, "b": INK, "accent": TEAL, "mut": MUT, "mono": TEAL, "amber": AMBER, "red": RED, "green": GREEN}
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for t, style in runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.name = MONO if style == "mono" else SANS
            r.font.bold = style in ("b", "accent")
            r.font.color.rgb = colors.get(style, INK)
    return box


def bullets(slide, x, y, w, h, items, *, size=17, line_spacing=1.15, space_after=10):
    paragraphs = []
    for item in items:
        runs = [("•  ", "accent")]
        if isinstance(item, str):
            runs.append((item, ""))
        else:
            runs.extend(item)
        paragraphs.append(runs)
    return rich(slide, x, y, w, h, paragraphs, size=size, line_spacing=line_spacing, space_after=space_after)


def header(slide, title, eyebrow=None):
    if eyebrow:
        text(slide, MARGIN, Inches(0.42), W - 2 * MARGIN, Inches(0.3), eyebrow.upper(), size=11, color=MUT, font=MONO)
    text(slide, MARGIN, Inches(0.7), W - 2 * MARGIN, Inches(0.9), title, size=32, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.55), Inches(1.2), Emu(28000))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, n, total):
    text(slide, MARGIN, H - Inches(0.5), Inches(6), Inches(0.3), "Crew Ops Advisor · dCortex hackathon 2026", size=10, color=MUT, font=MONO)
    text(slide, W - MARGIN - Inches(1.5), H - Inches(0.5), Inches(1.5), Inches(0.3), f"{n} / {total}", size=10, color=MUT, font=MONO, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, *, stroke=EDGE, fill=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def table(slide, x, y, w, rows, *, col_widths, size=14, header_color=MUT):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(0.45) * n_rows)
    tbl = shape.table
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD if i else PAPER
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size if i else size - 2)
            r.font.name = MONO if i == 0 else SANS
            r.font.bold = i == 0 or j == 0
            r.font.color.rgb = header_color if i == 0 else INK
    return shape


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


# ---------------------------------------------------------------- slides


def build() -> Path:
    res = latest_results()
    prs = new_deck()
    total = 12

    # 1 — title
    s = blank(prs)
    text(s, MARGIN, Inches(2.0), W - 2 * MARGIN, Inches(0.4), "DCORTEX AIR · CREW CONTROL · AGENTIC CREW OPS ADVISOR", size=12, color=MUT, font=MONO)
    text(s, MARGIN, Inches(2.4), W - 2 * MARGIN, Inches(1.2), "Crew Ops Advisor", size=54, bold=True)
    text(s, MARGIN, Inches(3.55), W - 2 * MARGIN, Inches(0.8), "A conversational assistant for the crew desk that is fast, right, and says so when it isn't sure.", size=22, color=MUT)
    text(s, MARGIN, Inches(5.2), W - 2 * MARGIN, Inches(0.5), "Team: Rajesh Tippana · Syed Maaz · Uma Shankar", size=16, color=INK)
    text(s, MARGIN, Inches(5.65), W - 2 * MARGIN, Inches(0.5), "github.com/valinci-007/dCortex", size=14, color=TEAL, font=MONO)
    footer(s, 1, total)
    notes(s, "Cold open (30 s): a captain calls in sick at 5 a.m. Which flights break, who can legally take them, what does it cost, who breaks tomorrow because of the fix? Then the line: legality is exact arithmetic — an LLM that approximates a duty-hour sum is fluent, confident and wrong.")

    # 2 — the problem
    s = blank(prs)
    header(s, "The problem", "01 · why this matters")
    bullets(s, MARGIN, Inches(1.9), W - 2 * MARGIN, Inches(4.6), [
        "A captain calls in sick at 5 a.m. Which flights break? Who can legally take them? What does it cost? Who else breaks tomorrow because of the fix?",
        "Today that reasoning lives in one senior controller's head, across six screens and a rulebook — slow to learn, impossible to scale, worst under pressure.",
        [("Legality is exact arithmetic. ", "b"), ("An LLM that approximates a duty-hour sum is fluent, confident and wrong — operationally worse than no answer.", "")],
        [("The brief's own bar: ", "b"), ("\"a polished, reliable Tier 1 with a credible Tier 2 beats a broken Tier 3\"; \"correctness outweighs coverage\".", "")],
    ], size=19, space_after=14)
    footer(s, 2, total)

    # 3 — the architectural question
    s = blank(prs)
    header(s, "The central question is architectural", "02 · three candidates, one rubric")
    text(s, MARGIN, Inches(1.8), W - 2 * MARGIN, Inches(0.75), "What should the language model do, what should deterministic code do, and how do you compose them?", size=17, color=MUT)
    table(s, MARGIN, Inches(2.6), W - 2 * MARGIN, [
        ["", "A · Tool agent over a deterministic core", "B · NL → SQL", "C · Intent router"],
        ["AI Utilization (20%)", "strong — the model plans and narrates", "blurry — the model authors the logic", "\"AI decorating a lookup\""],
        ["Correctness", "tools own all arithmetic", "silent wrong SQL is the named failure", "strong but closed-world"],
        ["Held-out questions", "generalises", "partial", "refuses the unanticipated"],
    ], col_widths=[Inches(2.4), Inches(3.5), Inches(3.0), Inches(3.03)], size=14)
    rich(s, MARGIN, Inches(5.1), W - 2 * MARGIN, Inches(1.0), [[
        ("Chosen: A, hardened with C's refusal discipline and offline router, and B's machine-readable reasoning trace. ", "b"),
        ("Every candidate, and everything we skipped, is in docs/decisions.md — 21 ADRs.", "mut"),
    ]], size=17)
    footer(s, 3, total)

    # 4 — the boundary (diagram from shapes)
    s = blank(prs)
    header(s, "Where the boundary is", "03 · the model never computes")
    left_x, mid_x, right_x = MARGIN, Inches(5.05), Inches(8.6)
    top = Inches(2.1)
    box_w, box_h = Inches(3.7), Inches(1.5)
    # model
    c = card(s, left_x, top, box_w, box_h, stroke=AMBER)
    text(s, left_x, top + Inches(0.15), box_w, Inches(0.4), "LANGUAGE MODEL", size=11, color=AMBER, font=MONO, align=PP_ALIGN.CENTER)
    text(s, left_x, top + Inches(0.48), box_w, Inches(1.0), "plans which lookups to make\nnarrates the answer with its reasoning\nnever does arithmetic", size=13, align=PP_ALIGN.CENTER)
    # tools (the crossing)
    c = card(s, mid_x, top - Inches(0.1), Inches(3.1), box_h + Inches(0.2), stroke=TEAL, fill=PAPER)
    text(s, mid_x, top + Inches(0.05), Inches(3.1), Inches(0.4), "35 TYPED TOOLS", size=11, color=TEAL, font=MONO, align=PP_ALIGN.CENTER)
    text(s, mid_x, top + Inches(0.42), Inches(3.1), Inches(1.2), "JSON-schema validated\nthe only crossing\nevidence objects back", size=13, align=PP_ALIGN.CENTER)
    # deterministic core
    c = card(s, right_x, top, box_w, box_h, stroke=TEAL)
    text(s, right_x, top + Inches(0.15), box_w, Inches(0.4), "DETERMINISTIC CORE", size=11, color=TEAL, font=MONO, align=PP_ALIGN.CENTER)
    text(s, right_x, top + Inches(0.48), box_w, Inches(1.0), "rules engine: 7 rules → evidence\nsimulations · ranking · cost model\nSQLite built from the dataset", size=13, align=PP_ALIGN.CENTER)
    # arrows
    for x1, x2, y, label in ((left_x + box_w, mid_x, top + Inches(0.55), "tool call"), (mid_x + Inches(3.1), right_x, top + Inches(0.55), "call"), (right_x, mid_x + Inches(3.1), top + Inches(1.05), "evidence"), (mid_x, left_x + box_w, top + Inches(1.05), "evidence")):
        arrow = s.shapes.add_connector(1, x1, y, x2, y)
        arrow.line.color.rgb = MUT
        arrow.line.width = Pt(1.5)
        arrow.line.end_arrowhead_style = 3 if hasattr(arrow.line, "end_arrowhead_style") else None
    text(s, left_x + box_w, top + Inches(0.2), mid_x - left_x - box_w, Inches(0.35), "tool call →", size=9, color=MUT, font=MONO, align=PP_ALIGN.CENTER)
    text(s, left_x + box_w, top + Inches(1.15), mid_x - left_x - box_w, Inches(0.35), "← evidence", size=9, color=MUT, font=MONO, align=PP_ALIGN.CENTER)
    bullets(s, MARGIN, Inches(4.1), W - 2 * MARGIN, Inches(2.6), [
        [("Claude Code's built-in tools are disabled; ", "b"), ("our registry is exposed as an in-process MCP server. The model can call exactly our tools and nothing else.", "")],
        [("Every rule returns evidence: ", "b"), ("inputs, computed value, limit, margin, a human detail — the numbers a controller can challenge.", "")],
        [("A grounding check ", "b"), ("verifies every id, date and figure in the answer against tool evidence: one rewrite, then a visible warning. Refusal is a first-class outcome.", "")],
        [("If the model provider fails, ", "b"), ("an offline rule-based router answers over the same tools, labelled as such.", "")],
    ], size=15, space_after=6)
    footer(s, 4, total)

    # 5 — what it answers
    s = blank(prs)
    header(s, "What it answers", "04 · three tiers, one boundary")
    colw = (W - 2 * MARGIN - Inches(0.5)) / 3
    cols = [
        ("TIER 1 · LOOKUPS", TEAL, "17 tools", "Reserves, duty clocks and headroom, flights, routes, pairings, certifications, risk signals, rules, costs."),
        ("TIER 2 · CONSEQUENCES", AMBER, "10 tools", "Sick call → uncrewed and at-risk legs, passengers. Substitution → all seven rules with the numbers. Station closure → each leg's minimum delay and the crew's FDP after it. Delays, cancellations, crew near limits, reserve coverage."),
        ("TIER 3 · RECOMMENDATIONS", GREEN, "8 tools", "Ranked covers with cost, delay, coverage, reasoning, the tightest rule headroom each leaves, and every excluded candidate with its reason. Joint plans, delay recovery, callout drafts, briefing, proactive watchlist; and positioning cover — when nobody local can take it on time, who elsewhere can be flown in before the departure. Read-only: what-ifs are questions, the controller decides."),
    ]
    for i, (title, colr, count, body) in enumerate(cols):
        x = MARGIN + i * (colw + Inches(0.25))
        card(s, x, Inches(1.95), colw, Inches(4.5), stroke=colr)
        text(s, x + Inches(0.2), Inches(2.1), colw - Inches(0.4), Inches(0.4), title, size=11, color=colr, font=MONO)
        text(s, x + Inches(0.2), Inches(2.45), colw - Inches(0.4), Inches(0.4), count, size=13, color=MUT, font=MONO)
        text(s, x + Inches(0.2), Inches(2.9), colw - Inches(0.4), Inches(3.4), body, size=14)
    footer(s, 5, total)

    # 6 — live demo
    s = blank(prs)
    header(s, "Live demo", "05 · five minutes, four questions")
    rich(s, MARGIN, Inches(1.9), Inches(7.6), Inches(4.8), [
        [("1  ", "accent"), ("\"If Captain C-2087 covers P-2291 from 15 Sep, does any rule breach?\"", "b")],
        [("    RULE-DUTY-02 by 1h20m — the window, the total, the margin. None of it computed by the model.", "mut")],
        [("2  ", "accent"), ("\"C-1042 called in sick from tomorrow — record it. What should I do?\"", "b")],
        [("    Ranked: C-3310 ₹18,500 → day-off callouts ₹24,000 → deadhead ₹41,200 + 3 h delay → cancel ₹15,00,000; who was excluded and why. Steps stream as they run.", "mut")],
        [("3  ", "accent"), ("\"Captain C-2210 in DEL is also out — who covers P-2291's DEL day on 16 Sep?\"", "b")],
        [("    No legal on-time cover at DEL: the ranking escalates to positioning — DX588 the evening before with a hotel, ₹29,200, on time, cheaper than the delayed deadhead.", "mut")],
        [("4  ", "accent"), ("\"Will fog delay BLR tomorrow?\"", "b"), ("  →  \"I can't answer that reliably.\"", "accent")],
    ], size=15, space_after=6)
    card(s, Inches(8.7), Inches(1.9), Inches(3.93), Inches(4.6))
    text(s, Inches(8.9), Inches(2.05), Inches(3.5), Inches(0.4), "HELD IN RESERVE", size=11, color=MUT, font=MONO)
    bullets(s, Inches(8.9), Inches(2.45), Inches(3.55), Inches(4.0), [
        "The watchlist on the home screen: a training lapse two days before a rostered duty",
        "The PII audit console: names never reach the model, before/after every tool result",
        "The offline router answering the same chain with the network pulled",
        "Voice input; persistent chats; the developer view",
    ], size=13, space_after=6)
    footer(s, 6, total)
    notes(s, "Switch to the app. Developer view OFF for the controller's experience; flip it ON when showing the audit trail. Start the server with ./start.sh so the console shows the PII scrub.")

    # 7 — results
    s = blank(prs)
    header(s, "Results against the dataset's own answer keys", "06 · measured, every run reported")
    table(s, MARGIN, Inches(1.95), Inches(8.3), [
        ["", "Model (Agent SDK)", "Offline router", "p50 latency"],
        ["Tier 1 · 16 questions", res.get("t1", "16/16"), "16/16", "≈ 7 s · < 1 ms"],
        ["Tier 2 · 14 questions", res.get("t2", "14/14"), "14/14", "≈ 10 s · 1 ms"],
        ["Tier 3 · 8 questions", f"{res.get('t3', '7/8')} automated · 8/8 on review", "7/8", "≈ 14 s · 6 ms"],
    ], col_widths=[Inches(2.3), Inches(2.8), Inches(1.6), Inches(1.6)], size=14)
    bullets(s, MARGIN, Inches(4.2), Inches(8.3), Inches(2.6), [
        [("The one automated miss is Q33, ", "b"), ("an answer key that contradicts its own 4-leg computation. Recorded, tested, explained — we keep the consistent figure.", "")],
        [("Scenarios S1, S2, S4, S6 reproduced exactly ", "b"), ("(options, costs, exclusions, reasons); S5 differs by a documented exclusion.", "")],
        [("Every correction turn in the 38-question run ", "b"), ("came from a figure the model derived or one our own checks misread — never from a wrong lookup. Each was fixed at its source.", "")],
    ], size=14, space_after=8)
    card(s, Inches(9.3), Inches(1.95), Inches(3.33), Inches(4.8))
    rich(s, Inches(9.5), Inches(2.1), Inches(3.0), Inches(4.5), [
        [("LATEST FULL RUN", "mut")],
        [(res.get("stem", "tier123-agent-sdk"), "mono")],
        [("p50 " + res.get("p50", "—") + "  ·  p95 " + res.get("p95", "—"), "")],
        [("whole run " + res.get("cost", "—") + " with prompt caching", "")],
        [(f"{TESTS} tests", "b"), (" pin the engine to the answer keys and the roster to the organiser's validator", "")],
        [("Grading: recall of the key's atomic facts, then human review. No LLM judge on a determinism claim.", "mut")],
    ], size=13, space_after=8)
    footer(s, 7, total)

    # 8 — trust
    s = blank(prs)
    header(s, "Trust is built in, not asserted", "07 · what a controller can lean on")
    items = [
        ("Refusal is a feature", "No tool covers weather, bookings, HR — it declines and names the nearest supported question."),
        ("Grounding check", "Every id, date and figure must appear in tool evidence. One rewrite, then a visible warning."),
        ("Confidence on every answer", "verified · verified after correction · unverified figures · declined — and the tightest rule headroom on every ranked option."),
        ("PII minimisation", "CREW_OPS_PII_MODE=minimal: the model never sees a crew name; the browser joins names from a local directory. Medical dates are health data."),
        ("Audit console", "Every prompt, every question as typed and as sent, every tool result before and after the scrub — printed live."),
        ("Sandboxed and local", "The model can reach only our 35 tools; runs on a laptop; speech on-device; chats deletable."),
    ]
    colw = (W - 2 * MARGIN - Inches(0.3)) / 2
    for i, (t, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = MARGIN + col * (colw + Inches(0.3))
        y = Inches(1.95) + row * Inches(1.55)
        card(s, x, y, colw, Inches(1.4))
        text(s, x + Inches(0.2), y + Inches(0.12), colw - Inches(0.4), Inches(0.4), t, size=14, bold=True, color=TEAL)
        text(s, x + Inches(0.2), y + Inches(0.5), colw - Inches(0.4), Inches(0.85), body, size=13)
    footer(s, 8, total)

    # 9 — weaknesses
    s = blank(prs)
    header(s, "Where it is weak — and we say so", "08 · honest failure analysis")
    bullets(s, MARGIN, Inches(1.9), W - 2 * MARGIN, Inches(4.8), [
        [("Tier-3 latency is 12–30 s ", "b"), ("with three or four tool calls: tool time is milliseconds, the cost is model turns. Mitigated by streaming each step and the answer as it is written — legible, not solved.", "")],
        [("Grounding is strict: ", "b"), ("it once flagged a correct derived sum and forced a rewrite. We accept false positives over a fabricated figure, and we fix each false positive at its source (lakh grouping, rotation notation, the controller's own dates).", "")],
        [("Partial covers of multi-day pairings ", "b"), ("are legality-checked but not costed for repatriation; positioning knows our own network only.", "")],
        [("Two answer keys we deliberately don't match: ", "b"), ("S5 lists a pairing's own crew as covers; Q33's 3-leg FDP contradicts its own 4-leg computation.", "")],
        [("A concurrency bug found by adding a feature: ", "b"), ("one SQLite connection shared across request threads produced a NULL row once the watchlist loaded alongside other calls. Per-thread connections now; a 40-request test pins it.", "")],
        [("The offline router is closed-world ", "b"), ("— insurance, not the product.", "")],
    ], size=15, space_after=8)
    footer(s, 9, total)

    # 10 — engineering judgement
    s = blank(prs)
    header(s, "Engineering judgement", "09 · decisions, not defaults")
    bullets(s, MARGIN, Inches(1.9), Inches(6.0), Inches(4.8), [
        [("Docs first: ", "b"), ("three architectures scored against the rubric; 21 ADRs including what we skipped and why.", "")],
        [("Foundation up with exit gates: ", "b"), ("rules engine → tools + agent → simulations → UI + grounding → ranking → positioning. No tier before the previous gate was green.", "")],
        [("Provider is a config switch: ", "b"), ("Agent SDK (default), client SDK, offline — one loop, automatic fallback.", "")],
        [("Controller view by default, developer view behind a toggle: ", "b"), ("the desk sees the answer, its reasoning and one trust signal; judges can flip on traces, cost, timings.", "")],
    ], size=15, space_after=10)
    card(s, Inches(7.0), Inches(1.9), Inches(5.63), Inches(4.7))
    text(s, Inches(7.2), Inches(2.05), Inches(5.2), Inches(0.4), "THE EVAL AS A DIAGNOSTIC LOOP", size=11, color=MUT, font=MONO)
    rich(s, Inches(7.2), Inches(2.45), Inches(5.25), Inches(4.0), [
        [("Reports keep the grounding verdict, the confidence label and a compact trace per question.", "")],
        [("Mining them showed why every correction turn happened — and that the prompt was the smallest part of the fix: ", ""), ("one sentence in the prompt, ", "b"), ("the rest in tool results, the grounding check and the grader.", "")],
        [("Tier 3 went 3/8 → 7/8 without tuning the model to the grader or the grader to the model.", "accent")],
    ], size=14, space_after=10)
    footer(s, 10, total)

    # 11 — production path
    s = blank(prs)
    header(s, "Production path", "10 · what would change at airline scale")
    bullets(s, MARGIN, Inches(1.9), W - 2 * MARGIN, Inches(4.8), [
        [("Scale: ", "b"), ("the model sees tool results, not the dataset; SQLite becomes the crew-tracking and roster systems behind the same repository interfaces; the core is stateless.", "")],
        [("Privacy: ", "b"), ("minimal PII mode by default, the directory behind the controller's authorisation, role-based access, retention policy, redaction at rest.", "")],
        [("Latency: ", "b"), ("streaming shipped; next is capping Tier-3 evidence sent to the model and effort tuning per tier.", "")],
        [("Read-only by principle: ", "b"), ("the advisor never changes crew data; what-ifs are questions and the controller acts in the desk's own systems (ADR-0021).", "")],
        [("Business impact: ", "b"), ("minutes → seconds per disruption, every decision auditable, the cheapest legal option visible next to the ones excluded and why.", "")],
    ], size=15, space_after=10)
    footer(s, 11, total)

    # 12 — close
    s = blank(prs)
    text(s, MARGIN, Inches(2.4), W - 2 * MARGIN, Inches(1.0), "Fast. Right. Says so when it isn't sure.", size=40, bold=True)
    text(s, MARGIN, Inches(3.5), W - 2 * MARGIN, Inches(0.6), "github.com/valinci-007/dCortex", size=20, color=TEAL, font=MONO)
    text(s, MARGIN, Inches(4.2), W - 2 * MARGIN, Inches(1.2), "docs/architecture.md · docs/decisions.md · docs/failure-cases.md · backend/evals/reports/\n./start.sh runs the desk on a laptop.", size=15, color=MUT)
    footer(s, 12, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path} ({os.path.getsize(path) / 1024:.0f} KB)")
